# -*- coding: utf-8 -*-
"""
Created on Thu Apr 23 16:48:20 2020

A script to create a powerpoint presentation for our display in the window

this file might not work if the layout of the website is changed too much

@author: skjerns
"""
import os
from joblib import Parallel, delayed
from pptx import Presentation
from tqdm import tqdm
import requests
import time
from bs4 import BeautifulSoup
from scipy import misc
import io
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
#%%

sortiment_url = 'https://www.buergerstiftung-karlsruhe.de/leihlokal/sortiment/?product-page='

def get(url, sleep=0.5):
    """retrieve an url and wait some second"""
    c = requests.get(url)
    if sleep>5: return c
    time.sleep(sleep)
    if not c.ok:
        print(f'Not OK: {c.reason}: {url}')
        return get(url, sleep*2)
    return c
    
def get_page_numbers():
    """retrieve the leihlokal sortiment and see how many pages there are"""
    c = get(sortiment_url)
    assert c.ok, f'Could not get url: {c.reason}'
    c = BeautifulSoup(c.content, 'html.parser')
    n_pages = (c.find_all('a', attrs={'class':'page-numbers'})[-2].text)
    return int(n_pages)

def jpg2int(jpg_bytes):
    return misc.imread(io.BytesIO(jpg_bytes))

# this is the url that we use to fetch the
n_pages = get_page_numbers()
request_urls = [sortiment_url + str(i) for i in range(1, n_pages+1)]

# we request 8 pages at once and then 200ms delay
res = Parallel(n_jobs=8, prefer='threads')(delayed(get)(url) for url in tqdm(request_urls))


# get all <li> tags that are of class 'product'
products = []
for page in res:
    soup = BeautifulSoup(page.content, 'html.parser')
    products += soup.find_all('li', attrs={'class':'product'})

# extract names, urls and image-urls
names = [p.find_all('h2')[0].text for p in products]
urls = [p.a.attrs['href'] for p in products]
images_urls = [p.a.img.attrs['srcset'].split(', ')[-1].split(' ')[0] for p in products]
images_urls = ['-'.join(url.split('-')[:-1])+url[-4:] if 'x' in url else url for url in images_urls]
codes = [p.find_all('a')[-1].attrs['data-product_sku'] for p in products]
    
# now download the images, store them.
# poor webserver, we download everything in batches of 100. should be quite fast.
images = Parallel(n_jobs=32, prefer='threads')(delayed(get)(url) for url in tqdm(images_urls))
images = [img.content for img in images]

os.makedirs('products', exist_ok=True)
for code, img in zip(codes, images):
    with open(os.path.join('products', f'{code}.jpg'), 'wb') as f:
        f.write(img)


products_mapping = {code:name for code, name in zip(codes, names)}


# Start creating the power point slides.
#%%

def make_slide(code):
   
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    top = height*0.15
    img_height = height*0.8
    left = (width-2)//2
    
    # the main image
    file = os.path.join('products', f'{code}.jpg')
    image = slide.shapes.add_picture(file, left, top, height=img_height)
    image.left = (prs.slide_width - image.width) // 2
    # image.top = (prs.slide_height - image.height) // 2
    
    # the BSKA-logo image
    
    file1 = os.path.join('logo', 'leihlokal.png')
    file2 = os.path.join('logo', 'qr.png')
    file3 = os.path.join('logo', 'stiftung.png')
    
    # logo1 = slide.shapes.add_picture(file1, width//66,  height//33, width = 0.9*((prs.slide_width - image.width) // 2))
    qrcode = slide.shapes.add_picture(file2, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    logo2 = slide.shapes.add_picture(file3, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    logo2.top =  prs.slide_height - logo2.height - height//20
    qrcode.top = image.top
    
    # the name of the item
    text_name = slide.shapes.add_textbox(0, 0, 0, 0)
    text_name.left = (prs.slide_width - text_name.width) // 2
    pr = text_name.text_frame.add_paragraph()
    pr.text = 'Leih dir doch ein(e/n)...'
    pr.font.size = Pt(24)
    pr.alignment = 2
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=0
    pr = text_name.text_frame.add_paragraph()
    pr.text = f'{products_mapping[code]} (#{code})'
    pr.font.size = Pt(50)
    pr.alignment = 2
    
    
    # # The info text on the right
    text_name = slide.shapes.add_textbox(0, image.top, 0, 0)
    text_name.left = (prs.slide_width//2 + image.width//2)+width//66
    text_name.top = image.top-Pt(21)
    pr = text_name.text_frame.add_paragraph()
    pr.text = 'Über 600 Gegenstände\naus Garten, Haushalt,\n' +\
              'Küche Kinder und \nHeimwerker\n\n' +\
              'Kein Mitgliedsbeitrag\nKeine Leihgebühr\nGegen Pfand ausleihen\n\n'+\
              'Unser Sortiment:'
              # 'Anmelden, Ausleihen, \nFreuen!'
    pr.font.size = Pt(19)
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=1.2
    pr = text_name.text_frame.add_paragraph()
    pr.text = 'bit.ly/leihlokal'
    pr.font.size = Pt(19)
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=1.2
    pr.font.underline=1
    pr.font.color.rgb = RGBColor(38, 136, 255)
    
    
    # Opening hours
    file = os.path.join('logo', 'oeffnungszeiten.png')
    hours = slide.shapes.add_picture(file, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    hours.left  = (prs.slide_width//2 + image.width//2)+width//66
    hours.top = image.top+image.height-hours.height
    
    

    
width = Cm(33)
height = Cm(20)
prs = Presentation()
prs.slide_width = width
prs.slide_height = height

for code in tqdm(codes[:200], desc='writing pptx'):
    make_slide(code)
prs.save('raspberry-pi-fenster.pptx')
    
os.startfile( os.path.abspath('raspberry-pi-fenster.pptx'))