# -*- coding: utf-8 -*-
"""
Created on Thu Apr 23 16:48:20 2020

A script to create a powerpoint presentation for our display in the window

this file might not work if the layout of the website is changed too much

@author: skjerns
"""
import os 
import ospath
from joblib import Parallel, delayed
from pptx import Presentation
from tqdm import tqdm
import pyexcel
import pandas as pd
import random
from scipy import misc
import io
import imageio
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
import PIL
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from joblib.memory import Memory

cache = Memory('z:/cache')

@cache.cache()
def load_resample(image_file, size=1080):
    image = PIL.Image.open(image_file)
    image.thumbnail([size, size], resample=PIL.Image.LANCZOS)
    image_stream = io.BytesIO()
    imageio.imwrite(image_stream, image, 'jpeg')
    return image_stream, image

#%%
def make_slide(prs, image_file, info):
    
    
    if not ('OK' in info['Kommentar Christian'] 
         or 'Super' in info['Kommentar Christian']):
        print (f'Skip foto {image_file}', info['Kommentar Christian'])
        return
    
    width = prs.slide_width
    height = prs.slide_height
    margin = int(height*0.025)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    


    image_stream, image = load_resample(image_file)

    portrait = True if image.height>image.width else False
    
    if portrait:
        image = slide.shapes.add_picture(image_stream, 0, 0, height=height-margin*4)
        image.top = margin*2
        image.left = height//2-image.width//2+margin*3
    else:
        image = slide.shapes.add_picture(image_stream, 0, 0, height=height-margin*4)
        image.top = height//2-image.height//2
        image.left = margin*2
    image_stream.close()
    
    frame = slide.shapes.add_textbox(image.left-margin, image.top-margin , 
                                     image.width+margin*2,  image.height+margin*2)
    frame.line.fill.solid()
    frame.fill.solid()
    # frame.fill.fore_color.rgb = RGBColor(255,255,255)
    frame.line.width = 0
    slide.shapes[0]._element.addprevious(slide.shapes[1]._element)
    
    # # The info text on the right
    
    textbox = slide.shapes.add_textbox(height, margin , width-height-margin*2, height//2-margin*2)
    pr = textbox.text_frame.paragraphs[0]
    
    
    title = info['titles']
    name = info['Name']
    alter = int(info['Alter'])
    hintergrund = info['Hintergrund']
    beschreibung = info['descriptions']
    
    
    textlen = len(info['descriptions']) + (60* (len(info['titles'])//35))
    text_scale_factor = 0.7

    if textlen>520:
        boxheight = Pt(10)*text_scale_factor*60
    elif textlen>460:
        boxheight = Pt(10)*text_scale_factor*57
    elif textlen>400:
        boxheight = Pt(10)*text_scale_factor*55
    elif textlen>300:
        boxheight = Pt(10)*text_scale_factor*50
    elif textlen>270:
        boxheight = Pt(10)*text_scale_factor*46
    elif textlen>180:
        boxheight = Pt(10)*text_scale_factor*42
    elif textlen>110:
        boxheight = Pt(10)*text_scale_factor*37
    elif textlen>56:
        boxheight = Pt(10)*text_scale_factor*32
    else:
        boxheight = Pt(10)*text_scale_factor*28
            
    textbox = slide.shapes.add_textbox(height+margin*6, height-boxheight-margin*2 , width-height-margin*7, boxheight)
    textbox.line.fill.solid()
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(255,255,255)
    textbox.line.width = 50
    
    textbox.text_frame.margin_left = int(width*0.01)
    textbox.text_frame.margin_top = int(width*0.01)
    textbox.text_frame.margin_right = int(width*0.01)
    textbox.text_frame.word_wrap=True
    
    
    
    # BSKE LOGO
 
    logo_box = slide.shapes.add_textbox(width-margin*10, margin, height=margin*9, width=margin*9)
    logo_box.line.fill.solid()
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = RGBColor(255,255,255)
    logo_box.line.width = 50
    logo = slide.shapes.add_picture('./logo/stiftung.png',logo_box.left+margin/2, logo_box.top+margin/2, height=logo_box.height-margin)
    poster = slide.shapes.add_picture('./logo/ausstellung_corona.png',logo_box.left,logo_box.top+logo_box.height+margin/2, width=logo_box.width)
    einreichung = slide.shapes.add_textbox(poster.left, poster.top+poster.height+margin/2, height=margin*2, width=poster.width)
    einreichung.line.fill.solid()
    einreichung.fill.solid()
    einreichung.fill.fore_color.rgb = RGBColor(255,255,255)
    pr = einreichung.text_frame.paragraphs[0]
    pr.text = 'Sie sehen hier: \nAlle Einreichungen.'
    pr.font.name = 'TT Norms Medium'
    pr.font.size = Pt(20)
    pr.alignment = PP_ALIGN.CENTER
    
    if info.Platz<=20:
        gewinner = slide.shapes.add_picture('./logo/gewinner.png',-margin*1.5, margin*6, height=logo_box.height-margin)
        gewinner.rotation=-45
    
    ## TITLE
    pr = textbox.text_frame.paragraphs[0]
    line = pr.add_run()
    line.text = title + '\n'
    line.font.name = 'TT Norms Medium' # install from the web if necessary
    line.font.size = Pt(45*text_scale_factor)
    # line.font.bold = True
    
    ## EMPTY LINE
    line = pr.add_run()
    line.text = '\n'
    line.font.name = 'TT Norms Medium' # install from the web if necessary
    line.font.size = Pt(10)
    # line.font.bold = True
    
    ## NAME + AGE
    line = pr.add_run()
    line.text = name
    line.font.name = 'TT Norms Medium' # install from the web if necessary
    line.font.size = Pt(30*text_scale_factor)
    line.line_spacing = 1.2

    line = pr.add_run()
    line.text = f', {alter}\n'
    line.font.name = 'TT Norms Medium' # install from the web if necessary
    line.font.size = Pt(30*text_scale_factor)
    line.line_spacing = 1.2    
    
    ## BACKGROUND
    line = pr.add_run()
    line.text = f'{hintergrund}\n'
    line.font.name = 'TT Norms Medium' # install from the web if necessary
    line.font.size = Pt(30*text_scale_factor)
    line.line_spacing = 0
    
    ## DESCRIPTION
    pr = textbox.text_frame.add_paragraph()
    pr.alignment = PP_ALIGN.JUSTIFY
    pr.text = beschreibung
    pr.font.name = 'TT Norms Regular' # install from the web if necessary
    pr.font.size = Pt(30*text_scale_factor)
    pr.line_spacing = 1  
        
    return prs
    
    


#%%
if __name__ == '__main__':
        
    images_dir = 'C:/Users/Simon/Nextcloud/Projekte/2020 Fotowettbewerb/Fotoeinreichungen/Fotoauswahl Jury'
    workbook = 'C:/Users/Simon/Nextcloud/Projekte/2020 Fotowettbewerb/Fotoeinreichungen/Teilnehmerliste Kontaktdaten (Uebersicht_Sitzung_v3).xls'
    descriptions = pd.read_csv(io.StringIO(pyexcel.load(workbook).get_csv()))[:-1]
    descriptions['Unnamed: 1'] = descriptions['Unnamed: 1'].astype(int)
    descriptions = descriptions.set_index('Unnamed: 1')
    images = ospath.list_files(images_dir, 'jpg')
    random.shuffle(images)
    # stop

    # approximately the display size
    width = Pt(192)*10
    height = Pt(108)*10
    
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    # fun = lambda x:len(descriptions.loc[int(ospath.basename(x)[:3])]['descriptions'])+(65 * (len(descriptions.loc[int(ospath.basename(x)[:3])]['titles'])//35))
    # images.sort(key=fun)
    
    for image_file in tqdm(images, desc='writing pptx'):
        try:
            id = int(ospath.basename(image_file)[:3])
            info = descriptions.loc[id]
            make_slide(prs, image_file, info)
            
        except Exception as e:
            print(f"failed to create slide for image {image_file}: {e}" )
        
    prs.save('fotoausstellung_slideshow.pptx')
    
    # now open it
    os.startfile( os.path.abspath('fotoausstellung_slideshow.pptx'))