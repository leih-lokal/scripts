# -*- coding: utf-8 -*-
"""
Created on Thu Apr 23 16:48:20 2020

A script to create a powerpoint presentation for our display in the window

this file might not work if the layout of the website is changed too much

@author: skjerns
"""
import os
import random
from leihlokal import LeihLokal
from joblib import Parallel, delayed
from pptx import Presentation
from tqdm import tqdm
import requests
import time
from scipy import misc
import io
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
#%%




#%%
if __name__ == '__main__':
        
    store = LeihLokal()
    
    items = [item for item in store.items.values() if item.status in ['instock', 'reserved', 'verliehen'] and item.image!='']
    # items = items[:10]
    random.shuffle(items)
    images_urls = [item.image for item in items]
    codes = [item.id for item in items]
    names = [item.name for item in items]
    
    #%%

    # approximately the display size
    width = Cm(21)
    height = Cm(29.7)
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    top = height*0.15
    img_height = height*0.8
    left = (width-2)//2
        
    # the leihlokal-logo 
    
    img_leihlokal = os.path.join('logo', 'leihlokal.png')
    
    logo = slide.shapes.add_picture(img_leihlokal, width*0.05,  height*0.025, 
                                    width = 0.9*((prs.slide_width) // 2))
    
    # the name of the item
    text_name = slide.shapes.add_textbox(0, 0, 0, 0)
    text_name.left = int(logo.left + logo.width + width*0.05)
    text_name.top = logo.top

    pr = text_name.text_frame.add_paragraph()
    pr.text = 'Statistik'
    pr.font.name = 'TT Norms Bold'
    pr.font.size = Pt(50)

    
    
#     # # The info text on the right
#     text_name = slide.shapes.add_textbox(0, image.top, 0, 0)
#     text_name.left = (prs.slide_width//2 + image.width//2)+width//66
#     text_name.top = image.top-Pt(21)
#     pr = text_name.text_frame.add_paragraph()
#     pr.text = f"""{len(codes)} Gegenstände für
# Garten, Küche, Kinder
# und zum Heimwerken.

# • Keine Leihgebühr
# • Gegen Pfand ausleihen
# • Spenden-finanizert
# • Ehrenamtlich organisiert

# Unser Sortiment online:"""
#               # 'Anmelden, Ausleihen, \nFreuen!'
#     pr.font.size = Pt(19)
#     pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
#     pr.line_spacing=1.2
#     pr = text_name.text_frame.add_paragraph()
#     pr.text = 'leihlokal-ka.de'
#     pr.font.size = Pt(19)
#     pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
#     pr.line_spacing=1.2
#     pr.font.underline=1
#     pr.font.color.rgb = RGBColor(38, 136, 255)
    
    
    # Opening hours
    file = os.path.join('logo', 'oeffnungszeiten.png')

            
    ppt_file = 'raspberry-pi-fenster.pptx'
    print(f'Saving to {ppt_file}')

    error = True
    while error:
        try:
            prs.save(ppt_file)
            error = False
        except:
            os.system("TASKKILL /F /IM powerpnt.exe")
            time.sleep(0.1)
    
    # now open it

    os.startfile( os.path.abspath('raspberry-pi-fenster.pptx'))