# -*- coding: utf-8 -*-
"""
Created on Thu Apr 23 16:48:20 2020

A script to create a powerpoint presentation for our display in the window

this file might not work if the layout of the website is changed too much

@author: skjerns
"""
import os
import random
from leihlokal import LeihLokal
from joblib import Parallel, delayed
from pptx import Presentation
from tqdm import tqdm
import requests
import time
from scipy import misc
import io
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io

#%%

sortiment_url = 'https://www.buergerstiftung-karlsruhe.de/leihlokal/sortiment/?product-page='

def download_image(url, code):
    file = os.path.join('products', f'{code}.jpg')
    if os.path.exists(file): return True
    c = get(url)
    image = Image.open(io.BytesIO(c.content))
    width, height = image.size
    factor = 800/width
    image = image.resize([int(width*factor), int(height*factor)])
    image.save(file, format='jpeg', quality=80)
    return True

def get(url, sleep=0.5):
    """retrieve an url and wait some second"""
    c = requests.get(url)
    if sleep>5: return c
    time.sleep(sleep)
    if not c.ok:
        print(f'Not OK: {c.reason}: {url}')
        return get(url, sleep*2)
    return c
    

def jpg2int(jpg_bytes):
    return misc.imread(io.BytesIO(jpg_bytes))


# Start creating the power point slides.
#%%

def make_slide(code):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    top = height*0.15
    img_height = height*0.8
    left = (width-2)//2
    
    # the main image
    file = os.path.join('products', f'{code}.jpg')
    image = slide.shapes.add_picture(file, left, top, height=img_height)
    image.left = (prs.slide_width - image.width) // 2
    # image.top = (prs.slide_height - image.height) // 2
    
    # the BSKA-logo image
    
    file1 = os.path.join('logo', 'leihlokal.png')
    file2 = os.path.join('logo', 'qr.png')
    file3 = os.path.join('logo', 'stiftung.png')
    
    # logo1 = slide.shapes.add_picture(file1, width//66,  height//33, width = 0.9*((prs.slide_width - image.width) // 2))
    qrcode = slide.shapes.add_picture(file2, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    logo2 = slide.shapes.add_picture(file3, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    logo2.top =  prs.slide_height - logo2.height - height//20
    qrcode.top = image.top
    
    # the name of the item
    text_name = slide.shapes.add_textbox(0, 0, 0, 0)
    text_name.left = (prs.slide_width - text_name.width) // 2
    pr = text_name.text_frame.add_paragraph()
    pr.text = 'Leih dir doch ein(e/n)...'
    pr.font.size = Pt(24)
    pr.alignment = 2
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=0
    pr = text_name.text_frame.add_paragraph()
    pr.text = f'{products_mapping[code]} (#{code})'
    pr.font.size = Pt(50)
    pr.alignment = 2
    
    
    # # The info text on the right
    text_name = slide.shapes.add_textbox(0, image.top, 0, 0)
    text_name.left = (prs.slide_width//2 + image.width//2)+width//66
    text_name.top = image.top-Pt(21)
    pr = text_name.text_frame.add_paragraph()
    pr.text = f"""{len(codes)} Gegenstände für
Garten, Küche, Kinder
und zum Heimwerken.

• Keine Leihgebühr
• Gegen Pfand ausleihen
• Spenden-finanizert
• Ehrenamtlich organisiert

Unser Sortiment online:"""
              # 'Anmelden, Ausleihen, \nFreuen!'
    pr.font.size = Pt(19)
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=1.2
    pr = text_name.text_frame.add_paragraph()
    pr.text = 'leihlokal-ka.de'
    pr.font.size = Pt(19)
    pr.font.name = 'TeXGyreAdventor' # install from the web if necessary
    pr.line_spacing=1.2
    pr.font.underline=1
    pr.font.color.rgb = RGBColor(38, 136, 255)
    
    
    # Opening hours
    file = os.path.join('logo', 'oeffnungszeiten.png')
    hours = slide.shapes.add_picture(file, width//66, 0, width = 0.9*((prs.slide_width - image.width) // 2))
    hours.left  = (prs.slide_width//2 + image.width//2)+width//66
    hours.top = image.top+image.height-hours.height


#%%
if __name__ == '__main__':
        
    store = LeihLokal()
    
    items = [item for item in store.items.values() if item.status in ['instock', 'reserved', 'verliehen'] and item.image!='']
    # items = items[:10]
    random.shuffle(items)
    images_urls = [item.image for item in items]
    codes = [item.id for item in items]
    names = [item.name for item in items]
    

    if not os.path.isdir('products'):
        os.makedirs('products')
        
    # now download the images, store them.
    # poor webserver, we download everything in batches of 100. should be quite fast.
    images = Parallel(n_jobs=8, prefer='threads')(
            delayed(download_image)(url, code) for url, code in tqdm(list(zip(images_urls, codes))))
    
    products_mapping = {code:name for code, name in zip(codes, names)}
    
    # approximately the display size
    width = Cm(33)
    height = Cm(20)
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    
    for code in tqdm(codes, desc='writing pptx'):
        try:
            # for all items: create a slide.
            make_slide(code)
        except Exception as e:
            print("failed to create slide for item " + str(code) + str(e))
            
    ppt_file = 'raspberry-pi-fenster.pptx'
    print(f'Saving to {ppt_file}')
    prs.save(ppt_file)
    
    # now open it
    os.startfile( os.path.abspath('raspberry-pi-fenster.pptx'))